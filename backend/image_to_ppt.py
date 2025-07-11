import cv2
import pytesseract
from pptx import Presentation
from pptx.util import Inches
import numpy as np
from pptx.enum.shapes import MSO_SHAPE


def extract_shapes(image):
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    blurred = cv2.GaussianBlur(gray, (5, 5), 0)
    edges = cv2.Canny(blurred, 50, 150)
    contours, _ = cv2.findContours(edges, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)

    shapes = []
    for cnt in contours:
        area = cv2.contourArea(cnt)
        if area < 100:
            continue
        approx = cv2.approxPolyDP(cnt, 0.02 * cv2.arcLength(cnt, True), True)
        x, y, w, h = cv2.boundingRect(approx)
        if len(approx) == 3:
            shape_type = 'triangle'
        elif len(approx) == 4:
            shape_type = 'rectangle'
        elif len(approx) > 6 and abs(1 - w / float(h)) < 0.1:
            shape_type = 'circle'
        else:
            shape_type = 'polygon'
        shapes.append({'type': shape_type, 'bbox': (x, y, w, h)})
    return shapes


def extract_text(image):
    data = pytesseract.image_to_data(image, output_type=pytesseract.Output.DICT)
    texts = []
    for i in range(len(data['text'])):
        text = data['text'][i].strip()
        if not text:
            continue
        x = data['left'][i]
        y = data['top'][i]
        w = data['width'][i]
        h = data['height'][i]
        texts.append({'text': text, 'bbox': (x, y, w, h)})
    return texts


def extract_tables(image):
    """Detect simple tables based on line patterns."""
    gray = cv2.cvtColor(image, cv2.COLOR_BGR2GRAY)
    edges = cv2.Canny(gray, 50, 150, apertureSize=3)
    lines = cv2.HoughLinesP(edges, 1, np.pi / 180, 80, minLineLength=40, maxLineGap=5)

    if lines is None:
        return []

    vertical = []
    horizontal = []
    for line in lines[:, 0]:
        x1, y1, x2, y2 = line
        if abs(x1 - x2) < 10:
            vertical.append(sorted((x1, x2)))
        elif abs(y1 - y2) < 10:
            horizontal.append(sorted((y1, y2)))

    if len(vertical) < 2 or len(horizontal) < 2:
        return []

    def cluster(vals, tol=10):
        vals = sorted([v[0] for v in vals])
        clusters = []
        for v in vals:
            if not clusters or v - clusters[-1] > tol:
                clusters.append(v)
        return clusters

    xs = cluster(vertical)
    ys = cluster(horizontal)
    if len(xs) < 2 or len(ys) < 2:
        return []

    x_min, x_max = xs[0], xs[-1]
    y_min, y_max = ys[0], ys[-1]
    rows = len(ys) - 1
    cols = len(xs) - 1
    return [{"bbox": (x_min, y_min, x_max - x_min, y_max - y_min),
            "rows": rows,
            "cols": cols}]


def add_elements_to_slide(slide, shapes, texts, tables, img_size):
    img_w, img_h = img_size
    slide_w = slide.part.slide_width
    slide_h = slide.part.slide_height

    def scale_x(x):
        return int(x * slide_w / img_w)

    def scale_y(y):
        return int(y * slide_h / img_h)

    for shp in shapes:
        x, y, w, h = shp['bbox']
        left, top = scale_x(x), scale_y(y)
        width, height = scale_x(w), scale_y(h)
        if shp['type'] == 'rectangle':
            slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
        elif shp['type'] == 'triangle':
            slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, left, top, width, height)
        elif shp['type'] == 'circle':
            slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, width, height)
        else:
            slide.shapes.add_shape(MSO_SHAPE.FREEFORM, left, top, width, height)

    for table in tables:
        x, y, w, h = table['bbox']
        left, top = scale_x(x), scale_y(y)
        width, height = scale_x(w), scale_y(h)
        rows = max(1, table['rows'])
        cols = max(1, table['cols'])
        slide.shapes.add_table(rows, cols, left, top, width, height)

    for tx in texts:
        x, y, w, h = tx['bbox']
        left, top = scale_x(x), scale_y(y)
        width, height = scale_x(w), scale_y(h)
        tb = slide.shapes.add_textbox(left, top, width, height)
        tb.text = tx['text']


def create_ppt_from_image(image_path, output_path='output.pptx'):
    image = cv2.imread(image_path)
    if image is None:
        raise FileNotFoundError(f"Cannot open {image_path}")
    shapes = extract_shapes(image)
    texts = extract_text(image)
    tables = extract_tables(image)

    prs = Presentation()
    # Force slide size to 1920x1080 (16:9) in pixels
    prs.slide_width = Inches(20)  # 1920 px at 96 DPI
    prs.slide_height = Inches(11.25)  # 1080 px at 96 DPI
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    img_h, img_w = image.shape[:2]
    add_elements_to_slide(slide, shapes, texts, tables, (img_w, img_h))
    prs.save(output_path)


if __name__ == '__main__':
    import argparse
    parser = argparse.ArgumentParser(description='Generate PPT slide from an image')
    parser.add_argument('image', help='Input image file')
    parser.add_argument('--output', '-o', default='output.pptx', help='Output PPTX file')
    args = parser.parse_args()
    create_ppt_from_image(args.image, args.output)

